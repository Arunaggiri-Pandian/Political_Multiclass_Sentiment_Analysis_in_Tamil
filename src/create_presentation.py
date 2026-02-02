"""
Create PowerPoint presentation for Political Multiclass Sentiment Analysis.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent.parent / "outputs"
FIGURES_DIR = OUTPUT_DIR / "figures"


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(100, 100, 100)


def add_content_slide(prs, title, content_lines):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(8)


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Table
    cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(num_rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.45 * num_rows)).table

    # Header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_image_slide(prs, title, image_path, caption=None):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(0.9), width=Inches(8))

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(100, 100, 100)


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Tamil Political Sentiment Analysis",
        "DravidianLangTech @ ACL 2026 | Team CHMOD_777")

    # Slide 2: Task Overview
    add_content_slide(prs, "Task Overview", [
        "Goal: Classify Tamil political tweets into 7 sentiment categories",
        "",
        "Classes: Opinionated, Sarcastic, Neutral, Positive,",
        "         Substantiated, Negative, None of the above",
        "",
        "Dataset: 4,352 train / 544 dev / 544 test samples",
        "",
        "Challenge: Severe class imbalance (7.96x ratio)",
        "",
        "Evaluation: Macro F1 Score"
    ])

    # Slide 3: Dataset Distribution
    add_image_slide(prs, "Original Dataset Distribution",
        FIGURES_DIR / "dataset_distribution.png",
        "Severe class imbalance: 7.96x ratio")

    # Slide 4: Data Augmentation
    add_content_slide(prs, "LLM-Based Data Augmentation", [
        "Challenge: Traditional EDA failed for Tamil",
        "",
        "Solution: Gemini 2.5 Flash for paraphrasing",
        "  • Preserves sentiment & hashtags",
        "  • Validates length & uniqueness",
        "",
        "Results:",
        "  • Original: 4,352 samples",
        "  • Augmented: 15,316 samples (3.52x)",
        "  • Imbalance reduced: 7.96x → 4.04x"
    ])

    # Slide 5: Augmentation Breakdown
    add_image_slide(prs, "Data Augmentation Breakdown",
        FIGURES_DIR / "augmentation_breakdown.png",
        "All classes now have 600+ samples")

    # Slide 6: Model Comparison Table
    add_table_slide(prs, "Model Comparison",
        ["Model", "Before Aug", "After Aug", "Improvement"],
        [
            ["MuRIL-base", "21.40%", "35.79% ✓", "+67%"],
            ["IndicBERT-v3-1B", "30.28%", "32.09%", "+6%"],
            ["IndicBERT-v3-270M", "-", "27.53%", "-"],
        ])

    # Slide 7: Augmentation Impact Chart
    add_image_slide(prs, "Impact of LLM Augmentation",
        FIGURES_DIR / "augmentation_impact.png",
        "MuRIL: +67% improvement | IndicBERT: +6% improvement")

    # Slide 8: Model Comparison Chart
    add_image_slide(prs, "Model Comparison",
        FIGURES_DIR / "model_comparison.png",
        "MuRIL (augmented) achieves best 35.79% Macro F1")

    # Slide 9: Per-Class Performance
    add_image_slide(prs, "Per-Class F1 Scores",
        FIGURES_DIR / "per_class_f1.png",
        "'None of above' easiest (86.5%), 'Substantiated' hardest (10.7%)")

    # Slide 10: Competition Comparison
    add_image_slide(prs, "Comparison with Last Year",
        FIGURES_DIR / "competition_comparison.png",
        "Our model would rank 3rd-4th place")

    # Slide 11: Key Insights
    add_content_slide(prs, "Key Insights", [
        "1. LLM augmentation provides 67% improvement",
        "   • Gemini paraphrasing >> traditional EDA",
        "",
        "2. Smaller models benefit more from augmentation",
        "   • MuRIL (236M) > IndicBERT-1B with augmented data",
        "",
        "3. Model size alone doesn't matter",
        "   • Quality of pre-training is key",
        "",
        "4. Would rank 3rd-4th in last year's competition",
        "   • 35.79% vs 1st place 37.7%"
    ])

    # Slide 12: Submissions
    add_content_slide(prs, "Final Submissions", [
        "Run 1: MuRIL (augmented)",
        "  • Macro F1: 35.79% (Best)",
        "",
        "Run 2: IndicBERT-v3-1B (augmented)",
        "  • Macro F1: 32.09%",
        "",
        "Run 3: IndicBERT-v3-270M (augmented)",
        "  • Macro F1: 27.53%"
    ])

    # Slide 13: Thank You
    add_title_slide(prs,
        "Thank You!",
        "Team CHMOD_777 | DravidianLangTech @ ACL 2026")

    # Save
    output_path = OUTPUT_DIR / "CHMOD_777_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    create_presentation()
